from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Configuration
OUTPUT_FILE = "AI_Cloud_Cost_Optimizer_Presentation.pptx"
IMAGE_PATH = r"C:\Users\Mounish Royal\.gemini\antigravity\brain\eda395ff-2c9d-4486-a05e-00348456908c\cloud_cost_architecture_1776581173359.png"

def add_slide(prs, title_text, content_points=None, layout_index=1):
    """Utility to add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Content
    if content_points:
        tf = slide.placeholders[1].text_frame
        tf.text = content_points[0]
        for point in content_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    return slide

def create_presentation():
    prs = Presentation()

    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI-Based Cloud Cost Optimizer"
    subtitle.text = "Final Year Project Viva Presentation\nPresented by: [Your Name]\nRoll No: [Your ID]\nGuided by: [Advisor Name]\nDepartment of Computer Science"

    # --- SLIDE 2: Introduction ---
    add_slide(prs, "Introduction", [
        "Rapid growth of Cloud Computing services (AWS, Azure, GCP).",
        "Organizations face 'Cloud Sprawl' and uncontrolled resource scaling.",
        "Cloud bill management is becoming a critical business challenge.",
        "Need for an automated, AI-driven solution to optimize spending without impacting performance."
    ])

    # --- SLIDE 3: Motivation ---
    add_slide(prs, "Motivation", [
        "Financial Impact: Businesses lose up to 30% of cloud spend on idle resources.",
        "Complexity: Multi-cloud pricing models are opaque and hard to track.",
        "Scalability: Manual rightsizing cannot keep up with dynamic workloads.",
        "Predictive Insights: Traditional tools are reactive, not proactive."
    ])

    # --- SLIDE 4: Problem Statement ---
    add_slide(prs, "Problem Statement", [
        "Identify and analyze the inefficiencies in cloud resource allocation.",
        "Inability of existing tools to provide human-readable explanations for cost changes.",
        "Developing a system that not only predicts costs but also optimizes them using rightsizing heuristics.",
        "Bridging the gap between raw data metrics and actionable cost-saving decisions."
    ])

    # --- SLIDE 5: Project Objectives ---
    add_slide(prs, "Project Objectives", [
        "Develop a Machine Learning model to predict cloud costs based on resource usage.",
        "Implement a 'Rightsizing' engine to suggest optimal resource configurations.",
        "Integrate Explainable AI (XAI) using Large Language Models (LLMs) for natural language reasoning.",
        "Build a user-friendly dashboard for real-time monitoring and optimization."
    ])

    # --- SLIDE 6: Scope of the Project ---
    add_slide(prs, "Scope of the Project", [
        "Focus Area: Compute (CPU), Memory (RAM), Storage, and Network I/O.",
        "Platforms: Generic cloud resource modeling (applicable to AWS/Azure/GCP).",
        "Target: SME's and Developers looking for cost transparency.",
        "Technologies: Python, Scikit-Learn, Flask, Ollama/Llama-3."
    ])

    # --- SLIDE 7: Literature Review ---
    add_slide(prs, "Literature Review", [
        "CloudHealth & Flexera: Enterprise-grade but high entry cost.",
        "AWS Cost Explorer: Platform-specific and misses third-party integrations.",
        "Research Papers: Focus on VM scheduling and load balancing but lack user-centric explainability.",
        "Our approach combines ML prediction + Dynamic Rightsizing + LLM explanations."
    ])

    # --- SLIDE 8: Gap Analysis ---
    add_slide(prs, "Gap Analysis", [
        "Lack of real-time predictive modeling in standard cloud dashboards.",
        "Absence of 'Why': Tools tell you cost went up, but not the narrative behind it.",
        "Dependency on static thresholds rather than dynamic workload patterns.",
        "The need for a localized AI solution (Ollama) to keep financial data private."
    ])

    # --- SLIDE 9: Proposed Methodology ---
    add_slide(prs, "Proposed Methodology", [
        "Phase 1: Synthetic/Real Workload Data Generation.",
        "Phase 2: Data Preprocessing & Feature Engineering.",
        "Phase 3: Training Regressive ML Models (Linear vs. Random Forest).",
        "Phase 4: Optimization Engine development (Rightsizing Heuristics).",
        "Phase 5: XAI Integration with Llama-3.",
        "Phase 6: Web Dashboard Integration."
    ])

    # --- SLIDE 10: System Architecture ---
    slide = add_slide(prs, "System Architecture")
    if os.path.exists(IMAGE_PATH):
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(7.0)
        slide.shapes.add_picture(IMAGE_PATH, left, top, width=width)
    else:
        p = slide.placeholders[1].text_frame.add_paragraph()
        p.text = "[Architecture Diagram Placeholder]"

    # --- SLIDE 11: Dataset and Attributes ---
    add_slide(prs, "Dataset and Attributes", [
        "Features (X): CPU Usage (%), Memory Usage (GB), Storage (GB), Network Traffic (GB).",
        "Target (y): Hourly Cloud Hosting Cost ($).",
        "Dataset Source: Simulated cloud provider logs covering diversos workload scenarios.",
        "Scaling: Standardized scaling for distance-based ML algorithms."
    ])

    # --- SLIDE 12: Data Preprocessing ---
    add_slide(prs, "Data Preprocessing", [
        "Cleaning: Handling missing values and removing noise.",
        "Normalization: Scaling features to [0,1] range for better model convergence.",
        "Splitting: 80% Training set, 20% Testing set.",
        "Validation: Using Mean Absolute Error (MAE) and R-squared metrics."
    ])

    # --- SLIDE 13: Machine Learning: Linear Regression ---
    add_slide(prs, "Machine Learning: Linear Regression", [
        "Utilized as a baseline model for project evaluation.",
        "Captures linear trends between resource growth and cost escalation.",
        "Pros: Interpretable, fast, low computational overhead.",
        "Cons: Struggles with non-linear pricing jumps (e.g., tiered storage costs)."
    ])

    # --- SLIDE 14: Machine Learning: Random Forest ---
    add_slide(prs, "Machine Learning: Random Forest", [
        "Ensemble learning method using multiple Decision Trees.",
        "Handles complex, non-linear relationships in cloud billing data.",
        "Reduces overfitting through bagging and feature randomness.",
        "Hyperparameters: 100 Estimators, Random State 42 for reproducibility."
    ])

    # --- SLIDE 15: Optimization Logic (Rightsizing) ---
    add_slide(prs, "Optimization Logic (Rightsizing)", [
        "Rightsizing: Matching instance size to actual performance requirements.",
        "Heuristic logic: Scale down resources by 20-30% if utilization is low (<50%).",
        "Lower Bounds: Ensuring system stability by maintaining minimum resource floors.",
        "Cost Savings: Predicted via the Random Forest model on 'optimized' features."
    ])

    # --- SLIDE 16: Explainable AI (XAI) with Ollama ---
    add_slide(prs, "Explainable AI (XAI) with Ollama", [
        "Powered by Llama-3 (localized installation via Ollama).",
        "Input: Original cost, Optimized cost, and Savings percentage.",
        "Output: 2-3 sentence executive summary explaining the ROI.",
        "Benefit: Increases user trust and transparency in AI recommendations."
    ])

    # --- SLIDE 17: Implementation Environment ---
    add_slide(prs, "Implementation Environment", [
        "Backend: Python 3.10+, Flask 3.0.",
        "Frontend: HTML5, CSS3, Jinja2 Templates.",
        "ML: Scikit-Learn, Pandas, NumPy.",
        "LLM Host: Ollama (running Llama-3 model).",
        "Monitoring: Python 'Logging' module for system health tracking."
    ])

    # --- SLIDE 18: Results - Model Performance ---
    add_slide(prs, "Results - Model Performance", [
        "Linear Regression MAE: $9.16",
        "Random Forest MAE: $6.64 (Chosen Model)",
        "Random Forest improved accuracy by ~27% over Linear Regression.",
        "The model accurately reflects the variance in cloud resource costs."
    ])

    # --- SLIDE 19: Results - Cost Savings ---
    add_slide(prs, "Results - Cost Savings", [
        "Average identified savings: 15% - 25% per environment.",
        "High Savings Case: Over-provisioned databases (>80% savings identified).",
        "Stable Performance: Optimized resources remain above safety thresholds.",
        "Real-time ROI calculation for stakeholders."
    ])

    # --- SLIDE 20: Software Demonstration ---
    add_slide(prs, "Software Demonstration", [
        "Key Screens: Input Dashboard, Prediction Graphs, Optimization Summary.",
        "API Integration: RESTful endpoints for external consumption.",
        "AI Chat: Interactive explanations generated on-the-fly.",
        "Responsiveness: Modern UI accessible from multiple devices."
    ])

    # --- SLIDE 21: Conclusion ---
    add_slide(prs, "Conclusion", [
        "Successfully developed an AI-driven framework for cloud cost reduction.",
        "Demonstrated that Random Forest is highly effective for cost prediction.",
        "LLM integration (Ollama) effectively bridges the gap between data and human understanding.",
        "Proposed system provides a blueprint for private, localized cloud optimization."
    ])

    # --- SLIDE 22: Future Scope & References ---
    add_slide(prs, "Future Scope & References", [
        "Multi-cloud Support: Ingesting AWS/Azure pricing SDKs directly.",
        "Auto-Scaling: Automated API triggers to shut down idle resources.",
        "Container Analysis: Kubernetes (K8s) resource optimization focus.",
        "References: [1] Cloud FinOps (O'Reilly), [2] Scikit-Learn Docs, [3] Ollama API Spec."
    ])

    # Save
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved as {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
